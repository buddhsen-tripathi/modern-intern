"""Document agent — generates presentations and documents via Gemini, then emails them.

Supports:
- "create presentation about X" → .pptx
- "create document about X" → .docx
"""

import io
import json
import logging
import os
import re
from email.mime.application import MIMEApplication
from email.mime.multipart import MIMEMultipart
from email.mime.text import MIMEText

import aiosmtplib
from google import genai
from google.genai import types
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from docx import Document
from docx.shared import Pt as DocxPt

from src.agents.base import BaseAgent
from src.config import BRAIN_MODEL

log = logging.getLogger(__name__)

PPTX_PROMPT = """\
Research the web and generate content for a professional presentation about: {topic}

Use current, factual data from search results. Return ONLY valid JSON (no markdown fences). Format:
{{
  "title": "Presentation Title",
  "slides": [
    {{
      "title": "Slide Title",
      "bullets": ["Point 1", "Point 2", "Point 3"]
    }}
  ]
}}

Rules:
- Exactly {size} slides (including title and conclusion)
- 3-5 bullet points per slide
- Keep bullets concise (under 15 words each)
- First slide is the title slide (no bullets needed, just title)
- Last slide is a summary/conclusion
"""

DOCX_PROMPT = """\
Research the web and generate content for a professional document about: {topic}

Use current, factual data from search results. Return ONLY valid JSON (no markdown fences). Format:
{{
  "title": "Document Title",
  "sections": [
    {{
      "heading": "Section Heading",
      "paragraphs": ["Paragraph text here.", "Another paragraph."]
    }}
  ]
}}

Rules:
- Exactly {size} sections
- 1-3 paragraphs per section, each 2-4 sentences
- Professional tone, clear and informative
"""


class DocumentAgent(BaseAgent):
    def __init__(self, api_key: str):
        self._client = genai.Client(api_key=api_key)
        self._gmail_address = os.getenv("GMAIL_ADDRESS", "")
        self._gmail_password = os.getenv("GMAIL_APP_PASSWORD", "")

    @property
    def name(self) -> str:
        return "Document"

    async def execute(self, params: dict, context: dict) -> dict:
        doc_type = params.get("type", "presentation")
        topic = params.get("topic", "").strip()
        size = params.get("size", 5)
        recipient = params.get("to", "").strip() or self._gmail_address

        if not topic:
            return {
                "status": "error",
                "error_type": "missing_topic",
                "message": "What should the document be about?",
            }

        if not self._gmail_address or not self._gmail_password:
            return {"status": "error", "message": "Gmail not configured. Can't send the file."}

        try:
            if doc_type == "presentation":
                content = await self._generate_content(PPTX_PROMPT.format(topic=topic, size=size))
                file_buf, filename = self._build_pptx(content)
                filetype = "presentation"
            else:
                content = await self._generate_content(DOCX_PROMPT.format(topic=topic, size=size))
                file_buf, filename = self._build_docx(content)
                filetype = "document"

            await self._send_email(recipient, topic, file_buf, filename, filetype)

            return {
                "status": "success",
                "message": f"{filetype.title()} about \"{topic}\" sent to {recipient}.",
                "filename": filename,
                "recipient": recipient,
                "topic": topic,
            }
        except Exception as e:
            log.error("Document creation failed: %s", e)
            return {"status": "error", "message": f"Failed to create {doc_type}: {e}"}

    async def _generate_content(self, prompt: str) -> dict:
        # Step 1: Research the topic with web search
        research_response = await self._client.aio.models.generate_content(
            model=BRAIN_MODEL,
            contents=f"Research this topic thoroughly and provide key facts, data, and insights:\n\n{prompt.split('Format:')[0]}",
            config=types.GenerateContentConfig(
                temperature=0.3,
                tools=[types.Tool(google_search=types.GoogleSearch())],
            ),
        )
        research = research_response.text.strip()

        # Step 2: Generate structured JSON from research (no search tool — clean output)
        json_prompt = (
            f"Based on this research:\n\n{research}\n\n"
            f"Now generate the following structured output.\n\n"
            f"{prompt}\n\n"
            f"CRITICAL: Return ONLY valid JSON. No markdown, no code fences, no extra text."
        )
        json_response = await self._client.aio.models.generate_content(
            model=BRAIN_MODEL,
            contents=json_prompt,
            config=types.GenerateContentConfig(
                temperature=0.2,
                response_mime_type="application/json",
            ),
        )
        text = json_response.text.strip()
        # Strip markdown code fences if somehow present
        text = re.sub(r"^```(?:json)?\s*", "", text)
        text = re.sub(r"\s*```$", "", text)
        return json.loads(text)

    def _build_pptx(self, content: dict) -> tuple[io.BytesIO, str]:
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        title = content.get("title", "Presentation")
        slides = content.get("slides", [])

        # Title slide
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
        txBox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(2))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(40)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0x1a, 0x1a, 0x2e)
        p.alignment = 1  # center

        # Content slides
        for s in slides:
            slide_title = s.get("title", "")
            bullets = s.get("bullets", [])
            if not bullets and slide_title == title:
                continue  # skip duplicate title

            slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

            # Title
            txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.5), Inches(1))
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = slide_title
            p.font.size = Pt(28)
            p.font.bold = True
            p.font.color.rgb = RGBColor(0x1a, 0x1a, 0x2e)

            # Bullets
            if bullets:
                txBox = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(11), Inches(5))
                tf = txBox.text_frame
                tf.word_wrap = True
                for i, bullet in enumerate(bullets):
                    if i == 0:
                        p = tf.paragraphs[0]
                    else:
                        p = tf.add_paragraph()
                    p.text = f"  •  {bullet}"
                    p.font.size = Pt(18)
                    p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
                    p.space_after = Pt(12)

        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)
        safe_name = re.sub(r"[^\w\s-]", "", title)[:40].strip().replace(" ", "_")
        return buf, f"{safe_name}.pptx"

    def _build_docx(self, content: dict) -> tuple[io.BytesIO, str]:
        doc = Document()
        title = content.get("title", "Document")
        sections = content.get("sections", [])

        doc.add_heading(title, level=0)

        for section in sections:
            heading = section.get("heading", "")
            paragraphs = section.get("paragraphs", [])

            if heading:
                doc.add_heading(heading, level=1)

            for para_text in paragraphs:
                p = doc.add_paragraph(para_text)
                for run in p.runs:
                    run.font.size = DocxPt(11)

        buf = io.BytesIO()
        doc.save(buf)
        buf.seek(0)
        safe_name = re.sub(r"[^\w\s-]", "", title)[:40].strip().replace(" ", "_")
        return buf, f"{safe_name}.docx"

    async def _send_email(self, to: str, topic: str, file_buf: io.BytesIO,
                          filename: str, filetype: str):
        msg = MIMEMultipart()
        msg["From"] = self._gmail_address
        msg["To"] = to
        msg["Subject"] = f"{filetype.title()}: {topic}"
        msg.attach(MIMEText(
            f"Hi,\n\nAttached is the {filetype} about \"{topic}\" "
            f"generated by Modern Intern.\n\nCheers",
            "plain",
        ))

        attachment = MIMEApplication(file_buf.read())
        attachment.add_header("Content-Disposition", "attachment", filename=filename)
        msg.attach(attachment)

        smtp = aiosmtplib.SMTP(hostname="smtp.gmail.com", port=587, start_tls=True)
        await smtp.connect()
        await smtp.login(self._gmail_address, self._gmail_password)
        await smtp.send_message(msg)
        await smtp.quit()
        log.info("Document emailed: %s to %s", filename, to)
